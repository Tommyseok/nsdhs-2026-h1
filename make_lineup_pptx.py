# -*- coding: utf-8 -*-
"""
내수동 경주자 2026 · 전체 라인업 PPT (16:9)
- 타이틀 1장 + 대가족반 6장 = 7장
- 가독성 최우선: 이름을 크게, 성별 색상(남=파랑/여=핑크), 담임/부담임 강조
출력: lineup.pptx
"""
import os, sys
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
sys.stdout.reconfigure(encoding='utf-8')

# ===== 데이터 =====
HOMEROOM = {1:'이윤정',2:'정명경',3:'전성희',4:'박대철',5:'이수지',6:'정유빈',
            7:'양지현',8:'이유성',9:'오덕현',10:'최희승',11:'석준원',12:'고은솔'}
SUB = {1:'이승헌',2:'조세경',3:'',4:'김희원',5:'박해니',6:'',
       7:'송진우',8:'이규리',9:'김주영',10:'',11:'',12:'오은규'}
LEADERS = {1:1,2:3,3:5,4:7,5:9,6:11}
CELL_STUDENTS = {
    1:[('권하율',1,'남'),('하은유',1,'여'),('홍기진',2,'남'),('서이수',3,'여'),('장지현',2,'여'),('김은',1,'여')],
    2:[('강은교',2,'여'),('명주원',3,'남'),('조인상',3,'여'),('연해준',2,'남'),('유나린',2,'여')],
    3:[('박지호',3,'여'),('한채은',3,'여'),('김준수',3,'남'),('김라일',2,'남'),('안민혁',1,'남')],
    4:[('원대연',2,'남'),('이준호',2,'남'),('전지훈',2,'남'),('오현우',1,'남'),('이빛가온',1,'남')],
    5:[('민예원',2,'여'),('김도연',2,'여'),('최지은',1,'여'),('김나현',2,'여'),('윤서희',2,'여')],
    6:[('조한나',1,'여'),('김한결',1,'남'),('이서후',1,'여'),('임태민',2,'남'),('박지웅',1,'남')],
    7:[('조강인',3,'남'),('이에녹',3,'남'),('최민호',2,'남'),('김정록',1,'남'),('정태율',1,'남')],
    8:[('심희은',3,'여'),('김은우',3,'여'),('최가은',3,'여'),('양소희',2,'여'),('오은수',1,'여')],
    9:[('나정인',3,'여'),('주성',3,'남'),('곽지원',3,'남'),('양인혁',3,'남'),('김다은',3,'여'),('박재영',1,'남')],
    10:[('이주형',1,'여'),('윤성하',1,'남'),('나호윤',1,'남'),('이소현',1,'여'),('배유나',2,'여')],
    11:[('최유현',1,'여'),('이가영',2,'여'),('최수지',1,'여'),('박동하',1,'남'),('송준원',1,'남')],
    12:[('정승빈',3,'여'),('김주하',2,'여'),('김채은',2,'여'),('유재희',2,'여'),('김강은',1,'여')],
}
SPECIAL = {
    1:[('고예나',1,'여'),('김민채',3,'여'),('송성모',3,'남'),('최승아',3,'여')],
    2:[('최지윤',3,'여'),('최서윤',3,'여'),('백결',3,'남'),('조혁준',3,'남')],
    3:[('김재원',1,'남'),('김유진',2,'여'),('김강희',1,'여'),('박서연',1,'여')],
    4:[('박시우',1,'남'),('하어림',1,'남'),('박진서',3,'여'),('임소명',3,'여')],
    5:[('오희수',3,'여'),('박민주',1,'여'),('장세민',3,'남'),('정원영',3,'남')],
    6:[('이서현',1,'여'),('이소민',3,'여'),('김도헌',3,'남'),('이하준',1,'남')],
}
FAM_CELLS = {1:(1,2),2:(3,4),3:(5,6),4:(7,8),5:(9,10),6:(11,12)}

# ===== 색상 =====
DARK   = RGBColor(0x17,0x13,0x2B)
PURPLE = RGBColor(0x7C,0x3A,0xED)
LILAC  = RGBColor(0xA7,0x8B,0xFA)
GREEN  = RGBColor(0x16,0xA3,0x74)
ORANGE = RGBColor(0xEA,0x58,0x0C)
INK    = RGBColor(0x1F,0x29,0x37)
MUTED  = RGBColor(0x6B,0x72,0x80)
MALE   = RGBColor(0x1D,0x4E,0xD8)
FEMALE = RGBColor(0xBE,0x18,0x5D)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
CARD_P = RGBColor(0xF5,0xF3,0xFF)
CARD_S = RGBColor(0xFF,0xF7,0xED)
LINEC  = RGBColor(0xE5,0xE7,0xEB)
FONT='맑은 고딕'; FONTB='맑은 고딕'

def sort_st(s): return sorted(s, key=lambda t:(-t[1],t[0]))

def no_autofit(tf): tf.word_wrap=True

def runs(p, parts):
    for txt,size,bold,color in parts:
        r=p.add_run(); r.text=txt; r.font.size=Pt(size); r.font.bold=bold
        r.font.color.rgb=color; r.font.name=FONT

def green_badge(slide, x, y, s):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, s, s)
    sh.fill.solid(); sh.fill.fore_color.rgb=GREEN; sh.line.fill.background()
    sh.adjustments[0]=0.28
    tf=sh.text_frame; tf.word_wrap=False; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text='R'; r.font.size=Pt(int(s/Emu(1)/914400*40)); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONTB
    sh.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    return sh

def cell_card(slide, x, y, w, h, title, count, homeroom, sub, is_leader, students, special=False):
    card=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb=CARD_S if special else CARD_P
    card.line.color.rgb=ORANGE if special else LILAC; card.line.width=Pt(1.25)
    card.adjustments[0]=0.045
    card.shadow.inherit=False
    tf=card.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    tf.margin_left=Inches(0.24); tf.margin_right=Inches(0.18); tf.margin_top=Inches(0.2); tf.margin_bottom=Inches(0.12)
    # 헤더 줄: 소그룹반 N  + 인원
    p0=tf.paragraphs[0]; p0.space_after=Pt(2); p0.alignment=PP_ALIGN.LEFT
    runs(p0, [(title, 17 if special else 21, True, ORANGE if special else PURPLE),
              (f'   {count}명', 13, False, MUTED)])
    # 담임/부담임 또는 안내
    if special:
        p1=tf.add_paragraph(); p1.space_after=Pt(8)
        runs(p1, [('대가족 선생님이 더 자주 함께해요', 12, False, MUTED)])
    else:
        p1=tf.add_paragraph(); p1.space_after=Pt(1)
        runs(p1, [('담임 ', 13, True, INK), (homeroom, 18, True, INK),
                  ('  ★가족장' if is_leader else '', 12, True, ORANGE)])
        p2=tf.add_paragraph(); p2.space_after=Pt(8)
        runs(p2, [('부담임 ', 13, True, MUTED), (sub if sub else '—', 17, True, INK)])
    # 학생 명단 — 이름 크게
    for (n,g,sx) in sort_st(students):
        sp=tf.add_paragraph(); sp.space_after=Pt(5)
        runs(sp, [(n, 22, True, INK),
                  (f'  {g}학년 ', 13, False, MUTED),
                  (sx, 14, True, MALE if sx=='남' else FEMALE)])

def add_title_slide(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb=DARK; bg.line.fill.background(); bg.shadow.inherit=False
    green_badge(s, Inches(0.9), Inches(2.55), Inches(1.5))
    tb=s.shapes.add_textbox(Inches(2.7), Inches(2.35), Inches(9.6), Inches(2.9)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; runs(p,[('내수동 경주자 ',46,True,WHITE),('2026',46,True,LILAC)])
    p2=tf.add_paragraph(); p2.space_before=Pt(2); runs(p2,[('전체 라인업',32,True,GREEN)])
    p3=tf.add_paragraph(); p3.space_before=Pt(14); runs(p3,[('6 대가족반 · 12 소그룹반 + 더 자주 보고 싶은 친구들 · 총 86명',18,False,RGBColor(0xCB,0xD5,0xE1))])
    p4=tf.add_paragraph(); p4.space_before=Pt(6)
    r=p4.add_run(); r.text='“나는 선한 싸움을 싸우고 나의 달려갈 길을 마치고” — 딤후 4:7'; r.font.size=Pt(13); r.font.italic=True; r.font.color.rgb=RGBColor(0x94,0xA3,0xB8); r.font.name=FONT
    ftb=s.shapes.add_textbox(Inches(0.9), Inches(6.9), Inches(11), Inches(0.4)); ftf=ftb.text_frame
    fr=ftf.paragraphs[0].add_run(); fr.text='내수동교회 고등부 · 발행 '+date.today().strftime('%Y-%m-%d'); fr.font.size=Pt(11); fr.font.color.rgb=RGBColor(0x6B,0x72,0x80); fr.font.name=FONT

def add_family_slide(prs, fam):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb=WHITE; bg.line.fill.background(); bg.shadow.inherit=False
    # 헤더
    green_badge(s, Inches(0.45), Inches(0.42), Inches(0.72))
    tb=s.shapes.add_textbox(Inches(1.35), Inches(0.34), Inches(8.5), Inches(1.0)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; runs(p,[(f'대가족반 {fam}',34,True,DARK)])
    c1,c2=FAM_CELLS[fam]; total=len(CELL_STUDENTS[c1])+len(CELL_STUDENTS[c2])+len(SPECIAL[fam])
    p2=tf.add_paragraph(); runs(p2,[(f'★ 가족장 {HOMEROOM[LEADERS[fam]]}    ·    총 {total}명 (소그룹 2 + 보고 싶은 친구들 1)',15,False,MUTED)])
    # 우상단 브랜드
    rb=s.shapes.add_textbox(Inches(9.8), Inches(0.5), Inches(3.1), Inches(0.5)); rtf=rb.text_frame
    rr=rtf.paragraphs[0]; rr.alignment=PP_ALIGN.RIGHT
    rn=rr.add_run(); rn.text='내수동 경주자 2026'; rn.font.size=Pt(13); rn.font.bold=True; rn.font.color.rgb=LILAC; rn.font.name=FONTB
    # 3 컬럼 카드
    left=Inches(0.45); top=Inches(1.55); gap=Inches(0.28)
    card_w=(prs.slide_width - left*2 - gap*2)/3
    card_h=prs.slide_height - top - Inches(0.4)
    cell_card(s, left, top, card_w, card_h, f'소그룹반 {c1}', len(CELL_STUDENTS[c1]), HOMEROOM[c1], SUB[c1], LEADERS[fam]==c1, CELL_STUDENTS[c1])
    cell_card(s, left+card_w+gap, top, card_w, card_h, f'소그룹반 {c2}', len(CELL_STUDENTS[c2]), HOMEROOM[c2], SUB[c2], LEADERS[fam]==c2, CELL_STUDENTS[c2])
    cell_card(s, left+(card_w+gap)*2, top, card_w, card_h, '더 자주 보고 싶은 친구들', len(SPECIAL[fam]), '', '', False, SPECIAL[fam], special=True)

def main():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    add_title_slide(prs)
    for fam in range(1,7): add_family_slide(prs, fam)
    out=os.path.join(os.path.dirname(os.path.abspath(__file__)),'lineup.pptx')
    prs.save(out)
    print('PPTX 생성 완료:', out, f'({os.path.getsize(out)/1024:.0f} KB, {len(prs.slides.__iter__.__self__._sldIdLst)}장)')

if __name__=='__main__':
    main()
